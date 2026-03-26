#!/usr/bin/env python3
"""
Generate ULTIMA THULE investor pitch deck (.pptx).

Design: Brutalist style matching the landing page (index.html).
- Black/white base with RED (#FF2D2D) accent
- Space Mono + Inter fonts
- Bold section labels like "01 / ПРОБЛЕМА"
- Thick borders, offset shadows, uppercase tracking
- High-contrast cards with card-number + card-label pattern
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Brand Colors (matching :root vars in index.html) ──
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xF8, 0xF8, 0xF8)
RED = RGBColor(0xFF, 0x2D, 0x2D)
GREEN = RGBColor(0x00, 0xFF, 0x66)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
BLUE = RGBColor(0x00, 0x66, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
BODY_TEXT = RGBColor(0x44, 0x44, 0x44)
SUBTLE = RGBColor(0x55, 0x55, 0x55)

# Slide dimensions (widescreen 16:9)
W = Inches(13.333)
H = Inches(7.5)

# Font names
MONO = "Courier New"  # fallback for Space Mono
HEADING = "Arial"     # fallback for Inter


# ════════════════════════════════════════════════════════════
# ANIMATION HELPERS
# ════════════════════════════════════════════════════════════

def add_entrance_anim(slide, shape, delay_ms=0, dur_ms=500):
    """Add fade entrance animation to a shape."""
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        timing = slide._element.makeelement(qn("p:timing"), {})
        slide._element.append(timing)

    tn_lst = timing.find(qn("p:tnLst"))
    if tn_lst is None:
        tn_lst = timing.makeelement(qn("p:tnLst"), {})
        timing.append(tn_lst)

    par_tn = tn_lst.find(qn("p:par"))
    if par_tn is None:
        par_tn = tn_lst.makeelement(qn("p:par"), {})
        tn_lst.append(par_tn)
        c_tn_root = par_tn.makeelement(qn("p:cTn"), {"id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot"})
        par_tn.append(c_tn_root)
        child_tn_lst = c_tn_root.makeelement(qn("p:childTnLst"), {})
        c_tn_root.append(child_tn_lst)
    else:
        c_tn_root = par_tn.find(qn("p:cTn"))
        child_tn_lst = c_tn_root.find(qn("p:childTnLst"))

    seq = child_tn_lst.find(qn("p:seq"))
    if seq is None:
        seq = child_tn_lst.makeelement(qn("p:seq"), {"concurrent": "1", "nextAc": "seek"})
        child_tn_lst.append(seq)
        seq_ctn = seq.makeelement(qn("p:cTn"), {"id": "2", "dur": "indefinite", "nodeType": "mainSeq"})
        seq.append(seq_ctn)
        seq_child = seq_ctn.makeelement(qn("p:childTnLst"), {})
        seq_ctn.append(seq_child)
        prev_cond = seq.makeelement(qn("p:prevCondLst"), {})
        seq.append(prev_cond)
        pc = prev_cond.makeelement(qn("p:cond"), {"evt": "onPrev", "delay": "0"})
        prev_cond.append(pc)
        tgt_el = pc.makeelement(qn("p:tgtEl"), {})
        pc.append(tgt_el)
        tgt_el.append(tgt_el.makeelement(qn("p:sldTgt"), {}))
        next_cond = seq.makeelement(qn("p:nextCondLst"), {})
        seq.append(next_cond)
        nc = next_cond.makeelement(qn("p:cond"), {"evt": "onNext", "delay": "0"})
        next_cond.append(nc)
        tgt_el2 = nc.makeelement(qn("p:tgtEl"), {})
        nc.append(tgt_el2)
        tgt_el2.append(tgt_el2.makeelement(qn("p:sldTgt"), {}))
    else:
        seq_ctn = seq.find(qn("p:cTn"))
        seq_child = seq_ctn.find(qn("p:childTnLst"))

    existing = seq_child.findall(qn("p:par"))
    next_id = 3 + len(existing) * 6

    anim_par = seq_child.makeelement(qn("p:par"), {})
    seq_child.append(anim_par)
    anim_ctn = anim_par.makeelement(qn("p:cTn"), {"id": str(next_id), "fill": "hold"})
    anim_par.append(anim_ctn)

    stCondLst = anim_ctn.makeelement(qn("p:stCondLst"), {})
    anim_ctn.append(stCondLst)
    cond = stCondLst.makeelement(qn("p:cond"), {"delay": "0"})
    stCondLst.append(cond)

    child2 = anim_ctn.makeelement(qn("p:childTnLst"), {})
    anim_ctn.append(child2)

    inner_par = child2.makeelement(qn("p:par"), {})
    child2.append(inner_par)
    inner_ctn = inner_par.makeelement(qn("p:cTn"), {
        "id": str(next_id + 1), "presetID": "10",
        "presetClass": "entr", "presetSubtype": "0",
        "fill": "hold", "nodeType": "afterEffect"
    })
    inner_par.append(inner_ctn)

    st2 = inner_ctn.makeelement(qn("p:stCondLst"), {})
    inner_ctn.append(st2)
    st2.append(st2.makeelement(qn("p:cond"), {"delay": str(delay_ms)}))

    child3 = inner_ctn.makeelement(qn("p:childTnLst"), {})
    inner_ctn.append(child3)

    # Set visibility
    s = child3.makeelement(qn("p:set"), {})
    child3.append(s)
    s_ctn = s.makeelement(qn("p:cBhvr"), {})
    s.append(s_ctn)
    s_ctn_inner = s_ctn.makeelement(qn("p:cTn"), {"id": str(next_id + 2), "dur": "1", "fill": "hold"})
    s_ctn.append(s_ctn_inner)
    s_st = s_ctn_inner.makeelement(qn("p:stCondLst"), {})
    s_ctn_inner.append(s_st)
    s_st.append(s_st.makeelement(qn("p:cond"), {"delay": "0"}))
    tgt_el = s_ctn.makeelement(qn("p:tgtEl"), {})
    s_ctn.append(tgt_el)
    sp_tgt = tgt_el.makeelement(qn("p:spTgt"), {"spid": str(shape.shape_id)})
    tgt_el.append(sp_tgt)
    to_el = s.makeelement(qn("p:to"), {})
    s.append(to_el)
    str_val = to_el.makeelement(qn("p:strVal"), {"val": "visible"})
    to_el.append(str_val)

    # AnimEffect (fade)
    ae = child3.makeelement(qn("p:animEffect"), {"transition": "in", "filter": "fade"})
    child3.append(ae)
    ae_cbhvr = ae.makeelement(qn("p:cBhvr"), {})
    ae.append(ae_cbhvr)
    ae_ctn = ae_cbhvr.makeelement(qn("p:cTn"), {"id": str(next_id + 3), "dur": str(dur_ms)})
    ae_cbhvr.append(ae_ctn)
    ae_tgt = ae_cbhvr.makeelement(qn("p:tgtEl"), {})
    ae_cbhvr.append(ae_tgt)
    ae_sp = ae_tgt.makeelement(qn("p:spTgt"), {"spid": str(shape.shape_id)})
    ae_tgt.append(ae_sp)


# ════════════════════════════════════════════════════════════
# SHAPE HELPERS
# ════════════════════════════════════════════════════════════

def set_bg(slide, r, g, b):
    """Set solid background color for a slide."""
    bg = slide._element.makeelement(qn("p:bg"), {})
    slide._element.insert(0, bg)
    bg_fill = bg.makeelement(qn("p:bgPr"), {})
    bg.append(bg_fill)
    solid = bg_fill.makeelement(qn("a:solidFill"), {})
    bg_fill.append(solid)
    srgb = solid.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (r, g, b)})
    solid.append(srgb)
    bg_fill.append(bg_fill.makeelement(qn("a:effectLst"), {}))


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font=HEADING, spacing=None,
             line_spacing=None):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    if spacing is not None:
        p.font._element.attrib[qn("a:spc")] = str(int(spacing))
    if line_spacing is not None:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        pPr.append(lnSpc)
        spcPts = lnSpc.makeelement(qn("a:spcPts"), {"val": str(int(line_spacing * 100))})
        lnSpc.append(spcPts)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, size=14,
                       color=GRAY, font=MONO, align=PP_ALIGN.LEFT, line_spacing=18):
    """Add text box with multiple lines (each line = separate paragraph)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        # line spacing
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        pPr.append(lnSpc)
        spcPts = lnSpc.makeelement(qn("a:spcPts"), {"val": str(int(line_spacing * 100))})
        lnSpc.append(spcPts)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None,
             border_width=3, corner_radius=0.0):
    """Add rectangle shape. corner_radius 0 = sharp, >0 = rounded."""
    if corner_radius > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, height, color, thickness=3):
    """Add a thin line (rectangle)."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def section_label(slide, text, y=Inches(0.5), color=RED, bg_color=BLACK, x=Inches(0.8)):
    """Add a section label like '01 / ПРОБЛЕМА' in the landing style."""
    # Background pill
    label_bg = add_rect(slide, x, y, Inches(2.6), Inches(0.35), bg_color)
    label_txt = add_text(slide, x + Inches(0.15), y + Inches(0.02),
                         Inches(2.3), Inches(0.35), text,
                         size=9, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.LEFT)
    return label_bg, label_txt


def brutal_card(slide, left, top, width, height, border_color=BLACK):
    """Create a brutalist card with thick border and offset shadow (light bg)."""
    # Shadow offset
    shadow = add_rect(slide, left + Inches(0.05), top + Inches(0.05),
                      width, height, LIGHT_GRAY, BLACK, border_width=3)
    # Main card
    card = add_rect(slide, left, top, width, height, WHITE, border_color, border_width=3)
    return card, shadow


def brutal_card_dark(slide, left, top, width, height, border_color=WHITE):
    """Dark card (like .brutal-card-dark on landing)."""
    card = add_rect(slide, left, top, width, height, DARK_GRAY, border_color, border_width=3)
    return card


def kpi_cell(slide, left, top, value, label, value_color=RED, delay=0):
    """KPI cell like on the landing page."""
    cell = add_rect(slide, left, top, Inches(2.7), Inches(1.5), WHITE, BLACK, border_width=2)
    v = add_text(slide, left + Inches(0.2), top + Inches(0.2),
                 Inches(2.3), Inches(0.8), value,
                 size=36, color=value_color, bold=True, font=HEADING)
    l = add_text(slide, left + Inches(0.2), top + Inches(0.95),
                 Inches(2.3), Inches(0.5), label,
                 size=8, color=GRAY, bold=True, font=MONO)
    add_entrance_anim(slide, cell, delay)
    add_entrance_anim(slide, v, delay + 150)
    add_entrance_anim(slide, l, delay + 250)
    return cell


def kpi_cell_dark(slide, left, top, value, label, value_color=RED, delay=0):
    """KPI cell on dark background."""
    cell = add_rect(slide, left, top, Inches(2.7), Inches(1.5), DARK_GRAY, WHITE, border_width=2)
    v = add_text(slide, left + Inches(0.2), top + Inches(0.2),
                 Inches(2.3), Inches(0.8), value,
                 size=36, color=value_color, bold=True, font=HEADING)
    l = add_text(slide, left + Inches(0.2), top + Inches(0.95),
                 Inches(2.3), Inches(0.5), label,
                 size=8, color=GRAY, bold=True, font=MONO)
    add_entrance_anim(slide, cell, delay)
    add_entrance_anim(slide, v, delay + 150)
    add_entrance_anim(slide, l, delay + 250)
    return cell


def progress_bar(slide, left, top, width, pct, color, label="", value_text="", delay=0):
    """Metric bar like on the landing page."""
    if label:
        add_text(slide, left, top - Inches(0.3), Inches(3), Inches(0.3),
                 label, size=9, color=BLACK, bold=True, font=MONO)
    # Track
    track = add_rect(slide, left, top, width, Inches(0.25), LIGHT_GRAY, BLACK, border_width=2)
    # Fill
    fill_w = int(width * pct) if int(width * pct) > 0 else Inches(0.1)
    fill = add_rect(slide, left, top, fill_w, Inches(0.25), color)
    fill.line.fill.background()
    add_entrance_anim(slide, fill, delay, dur_ms=800)
    if value_text:
        add_text(slide, left + width + Inches(0.15), top - Inches(0.02),
                 Inches(1.2), Inches(0.3), value_text,
                 size=14, color=color, bold=True, font=HEADING)
    return fill


# ════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank layout


# ────────────────────────────────────────────────────────────
# SLIDE 1: HERO (like #hero on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0x0A, 0x0A, 0x0A)

# Red divider line (like border-bottom: 6px solid RED)
red_line = add_line(s, 0, Inches(4.2), W, Pt(6), RED, thickness=6)

# Tag label (like .hero-tag)
tag_border = add_rect(s, Inches(1), Inches(1.6), Inches(3.8), Inches(0.38),
                       BLACK, RED, border_width=2)
tag_text = add_text(s, Inches(1.1), Inches(1.62), Inches(3.6), Inches(0.35),
                    "AI-POWERED TRAVEL PLANNING", size=9, color=RED, bold=True, font=MONO)

# Title (like .hero-title)
t_ultima = add_text(s, Inches(1), Inches(2.3), Inches(10), Inches(1.2),
                    "ULTIMA", size=80, color=WHITE, bold=True, font=HEADING)
t_thule = add_text(s, Inches(1), Inches(3.2), Inches(10), Inches(1.0),
                   "THULE", size=80, color=RED, bold=True, font=HEADING)

# Subtitle
t_sub = add_text(s, Inches(1), Inches(4.6), Inches(7), Inches(0.8),
                 "iOS-приложение для планирования путешествий с AI,\nтрекингом рейсов и офлайн-доступом.",
                 size=14, color=GRAY, font=MONO)

# Hero metrics row (like .hero-metrics)
metrics_line = add_line(s, Inches(1), Inches(5.7), Inches(11), Pt(1), GRAY, thickness=1)

metrics_data = [
    ("$14.5B", "TAM (ГЛОБАЛЬНО)"),
    ("35+", "ФИЧЕЙ ГОТОВО"),
    ("8x", "ROI (ОПТИМ.)"),
    ("1,500", "ДО BREAKEVEN"),
]
for i, (val, lbl) in enumerate(metrics_data):
    x = Inches(1 + i * 2.9)
    mv = add_text(s, x, Inches(5.9), Inches(2.5), Inches(0.6),
                  val, size=32, color=RED, bold=True, font=HEADING)
    ml = add_text(s, x, Inches(6.45), Inches(2.5), Inches(0.4),
                  lbl, size=8, color=GRAY, bold=True, font=MONO)
    add_entrance_anim(s, mv, 600 + i * 150)
    add_entrance_anim(s, ml, 700 + i * 150)

add_entrance_anim(s, tag_border, 0, 400)
add_entrance_anim(s, tag_text, 100)
add_entrance_anim(s, t_ultima, 200, 600)
add_entrance_anim(s, t_thule, 350, 600)
add_entrance_anim(s, t_sub, 500)


# ────────────────────────────────────────────────────────────
# SLIDE 2: PROBLEM (like #problem on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xF8, 0xF8, 0xF8)

section_label(s, "01 / ПРОБЛЕМА")
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(1.2),
                 "ПЛАНИРОВАНИЕ ПОЕЗДКИ —\nЭТО ХАОС", size=42, color=BLACK, bold=True, font=HEADING)
add_entrance_anim(s, title, 0, 600)

# 3 brutal cards (like .brutal-card on landing)
cards_data = [
    ("6+", "ПРИЛОЖЕНИЙ НА\nПОЕЗДКУ", "Карты, билеты, отели,\nзаметки, бюджет — всё\nразбросано"),
    ("5–10ч", "ВРЕМЯ НА ПЛАН\n(КОРОТКИЕ)", "На планирование одной\nпоездки до 10 дней.\nДлинные — до 30ч"),
    ("0", "AI-ПОМОЩНИКОВ\nНА РУССКОМ", "Ни одного travel planner\nна русском с AI и\nперсонализацией"),
]

for i, (num, lbl, desc) in enumerate(cards_data):
    x = Inches(0.8 + i * 4.0)
    card, shadow = brutal_card(s, x, Inches(3.3), Inches(3.7), Inches(3.2))
    add_text(s, x + Inches(0.25), Inches(3.5), Inches(3.2), Inches(0.9),
             num, size=48, color=RED, bold=True, font=HEADING)
    add_text(s, x + Inches(0.25), Inches(4.3), Inches(3.2), Inches(0.6),
             lbl, size=9, color=GRAY, bold=True, font=MONO)
    add_text(s, x + Inches(0.25), Inches(5.0), Inches(3.2), Inches(1.2),
             desc, size=12, color=BODY_TEXT, font=MONO)
    add_entrance_anim(s, shadow, 300 + i * 200)
    add_entrance_anim(s, card, 300 + i * 200)


# ────────────────────────────────────────────────────────────
# SLIDE 3: SOLUTION (like #solution on landing — dark bg)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0x0A, 0x0A, 0x0A)

section_label(s, "02 / РЕШЕНИЕ", bg_color=RED)
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(1.2),
                 "ВЕСЬ ТРИП —\nВ ОДНОМ ПРИЛОЖЕНИИ", size=42, color=WHITE, bold=True, font=HEADING)
add_entrance_anim(s, title, 0, 600)

sol_data = [
    ("01", "ПЛАНИРОВАНИЕ", "Мульти-трип. День за днём:\nместа, события, расходы.\nDrag-and-drop. Bucket list."),
    ("02", "AI-МОЗГ", "AI-генератор маршрутов.\nПерсональные рекомендации.\nOCR + парсинг почты."),
    ("03", "В ПОЕЗДКЕ", "Live Activity. Погода.\nВалюта. AR-навигация.\nОфлайн-карты. Журнал."),
]

for i, (num, name, desc) in enumerate(sol_data):
    x = Inches(0.8 + i * 4.0)
    card = brutal_card_dark(s, x, Inches(3.3), Inches(3.7), Inches(3.2))
    add_text(s, x + Inches(0.25), Inches(3.5), Inches(3.2), Inches(0.5),
             num, size=28, color=WHITE, bold=True, font=HEADING)
    add_text(s, x + Inches(0.25), Inches(4.1), Inches(3.2), Inches(0.4),
             name, size=14, color=WHITE, bold=True, font=HEADING)
    add_text(s, x + Inches(0.25), Inches(4.7), Inches(3.2), Inches(1.5),
             desc, size=11, color=GRAY, font=MONO)
    add_entrance_anim(s, card, 300 + i * 200)

# Feature tags row (like the tag pills on landing)
features_short = [
    "Мульти-трип", "AI Trip Generator", "Flight tracking",
    "OCR бронирований", "Aviasales", "Offline", "AR навигация",
    "Dynamic Island", "Supabase Sync", "Email Scanner",
]
for i, feat in enumerate(features_short):
    col = i % 5
    row = i // 5
    x = Inches(0.8 + col * 2.4)
    y = Inches(6.7 + row * 0.35)
    tag_rect = add_rect(s, x, y, Inches(2.2), Inches(0.28), BLACK, RED, border_width=2)
    add_text(s, x + Inches(0.1), y + Inches(0.01), Inches(2.0), Inches(0.26),
             feat, size=8, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_entrance_anim(s, tag_rect, 900 + i * 60)


# ────────────────────────────────────────────────────────────
# SLIDE 4: MARKET (like #market on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xF8, 0xF8, 0xF8)

section_label(s, "03 / РЫНОК")
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(1.0),
                 "РАЗМЕР РЫНКА TRAVEL APPS", size=42, color=BLACK, bold=True, font=HEADING)

# Divider
div_line = add_line(s, Inches(0.8), Inches(2.3), Inches(11.5), Pt(3), BLACK, thickness=3)
div_text = add_text(s, Inches(5), Inches(2.15), Inches(3), Inches(0.3),
                    "ГЛОБАЛЬНЫЙ РЫНОК", size=8, color=GRAY, bold=True, font=MONO, align=PP_ALIGN.CENTER)

# KPI cells row (like .kpi-cell on landing)
kpi_cell(s, Inches(0.8), Inches(2.7), "1.4B", "PAM — МЕЖДУНАРОДНЫЕ\nПУТЕШЕСТВЕННИКИ", BLACK, 200)
kpi_cell(s, Inches(3.7), Inches(2.7), "$14.5B", "TAM — TRAVEL\nPLANNING APPS", RED, 400)
kpi_cell(s, Inches(6.6), Inches(2.7), "$3.2B", "SAM — ITINERARY\nСЕГМЕНТ", BLACK, 600)
kpi_cell(s, Inches(9.5), Inches(2.7), "$1.55M", "SOM 5 ЛЕТ —\n50K ПЛАТЯЩИХ", RED, 800)

# Second divider
div2 = add_line(s, Inches(0.8), Inches(4.6), Inches(11.5), Pt(3), BLACK, thickness=3)
d2_text = add_text(s, Inches(5), Inches(4.45), Inches(3), Inches(0.3),
                   "РОССИЯ + СНГ", size=8, color=GRAY, bold=True, font=MONO, align=PP_ALIGN.CENTER)

kpi_cell(s, Inches(0.8), Inches(5.0), "42M", "РОССИЯН\nПУТЕШЕСТВУЮТ В ГОД", BLACK, 1000)
kpi_cell(s, Inches(3.7), Inches(5.0), "10M", "SAM — TECH-SAVVY\nПЛАНИРОВЩИКИ", RED, 1200)
kpi_cell(s, Inches(6.6), Inches(5.0), "$9.3M", "SAM ПО ВЫРУЧКЕ\n(3% × $31/ГОД)", BLACK, 1400)
kpi_cell_dark(s, Inches(9.5), Inches(5.0), "18%", "CAGR РОСТА\nРЫНКА", GREEN, 1600)


# ────────────────────────────────────────────────────────────
# SLIDE 5: PRODUCT — Features Grid (like #product on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0x0A, 0x0A, 0x0A)

section_label(s, "04 / ПРОДУКТ", bg_color=RED)
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "35+ ФИЧЕЙ ГОТОВЫ", size=42, color=WHITE, bold=True, font=HEADING)
add_entrance_anim(s, title, 0)

features = [
    "Мульти-трип менеджмент", "AI рекомендации (Gemini)", "AI Trip Generator",
    "Flight tracking + OCR", "Погода + Валюта", "Packing list + AI",
    "Offline-режим", "Live Activities", "Supabase Auth + Sync",
    "Bucket list + POI", "AR навигация", "Email Scanner",
    "Travel Stats", "Dual Timezone", "Today's Schedule",
    "Travel Journal", "Drag & Drop", "Smart Geofence",
]

for i, feat in enumerate(features):
    col = i % 6
    row = i // 6
    x = Inches(0.5 + col * 2.08)
    y = Inches(2.3 + row * 1.5)
    cell = add_rect(s, x, y, Inches(1.95), Inches(1.25), DARK_GRAY, WHITE, border_width=2)
    add_text(s, x + Inches(0.12), y + Inches(0.3), Inches(1.7), Inches(0.7),
             feat, size=11, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_entrance_anim(s, cell, 100 + i * 60)


# ────────────────────────────────────────────────────────────
# SLIDE 6: TECH STACK (like #stack on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xF8, 0xF8, 0xF8)

section_label(s, "05 / СТЕК")
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "АРХИТЕКТУРА В ЦИФРАХ", size=42, color=BLACK, bold=True, font=HEADING)

# Stats in kpi cells — dark cells for bright colors on light bg
kpi_cell_dark(s, Inches(0.8), Inches(2.3), "27+", "СЕССИЙ\nРАЗРАБОТКИ", GREEN, 200)
kpi_cell(s, Inches(3.7), Inches(2.3), "35+", "ФИЧЕЙ\nГОТОВО", RED, 400)
kpi_cell_dark(s, Inches(6.6), Inches(2.3), "12", "ТАБЛИЦ БД\n+ RLS", YELLOW, 600)
kpi_cell_dark(s, Inches(9.5), Inches(2.3), "8", "API ЧЕРЕЗ\nПРОКСИ", BLUE, 800)

# Stack items (like .stack-item on landing)
stack = [
    ("FRONTEND", "SwiftUI", "iOS 26, Glassmorphism UI, 6 палитр"),
    ("STORAGE", "SwiftData", "Offline-first, CoreData replacement"),
    ("BACKEND", "Supabase", "Auth + DB + Storage + Edge Functions"),
    ("AI", "Gemini 2.5 Flash", "Через облачный прокси, работает из РФ"),
    ("INTEGRATIONS", "Travelpayouts", "Aviasales рейсы + Hotellook отели"),
]

for i, (layer, tech, desc) in enumerate(stack):
    y = Inches(4.2 + i * 0.58)
    row_bg = add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.5), WHITE, BLACK, border_width=2)
    add_text(s, Inches(1.0), y + Inches(0.08), Inches(1.8), Inches(0.35),
             layer, size=8, color=GRAY, bold=True, font=MONO)
    add_text(s, Inches(3.0), y + Inches(0.06), Inches(2.5), Inches(0.35),
             tech, size=13, color=BLACK, bold=True, font=HEADING)
    add_text(s, Inches(6.0), y + Inches(0.1), Inches(6), Inches(0.35),
             desc, size=10, color=SUBTLE, font=MONO)
    add_entrance_anim(s, row_bg, 900 + i * 120)


# ────────────────────────────────────────────────────────────
# SLIDE 7: BUSINESS MODEL (like pricing cards on landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0x0A, 0x0A, 0x0A)

section_label(s, "06 / МОНЕТИЗАЦИЯ", bg_color=RED)
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "FREEMIUM B2C + B2B", size=42, color=WHITE, bold=True, font=HEADING)

tiers = [
    ("FREE", "0₽", ["3 поездки", "Базовый AI", "Онлайн-режим"], GRAY, False),
    ("PREMIUM", "349₽/мес", ["Безлимит поездок", "Полный AI + OCR", "Офлайн + Sync"], RED, True),
    ("ГОДОВОЙ", "2 490₽/год", ["Всё из Premium", "Скидка 40%", "Приоритетная поддержка"], GREEN, False),
    ("B2B", "per-seat", ["Корп. режим", "Админ-панель", "Отчёты + SSO"], BLUE, False),
]

for i, (name, price, features_list, clr, featured) in enumerate(tiers):
    x = Inches(0.6 + i * 3.1)
    card_w = Inches(2.9)

    # Card body
    card = add_rect(s, x, Inches(2.5), card_w, Inches(4.2), DARK_GRAY, clr, border_width=3)

    # Header bar (like .pricing-header)
    header_color = RED if featured else BLACK
    header = add_rect(s, x, Inches(2.5), card_w, Inches(1.4), header_color, clr, border_width=3)

    add_text(s, x + Inches(0.2), Inches(2.65), card_w - Inches(0.4), Inches(0.3),
             name, size=9, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.0), card_w - Inches(0.4), Inches(0.7),
             price, size=30, color=WHITE, bold=True, font=HEADING, align=PP_ALIGN.CENTER)

    # Features
    for j, feat in enumerate(features_list):
        fy = Inches(4.15 + j * 0.55)
        # separator line
        sep = add_line(s, x + Inches(0.2), fy, card_w - Inches(0.4), Pt(1), RGBColor(0x33, 0x33, 0x33))
        add_text(s, x + Inches(0.2), fy + Inches(0.05), card_w - Inches(0.4), Inches(0.4),
                 feat, size=11, color=GRAY, font=MONO, align=PP_ALIGN.CENTER)

    add_entrance_anim(s, card, 200 + i * 200)
    add_entrance_anim(s, header, 200 + i * 200)


# ────────────────────────────────────────────────────────────
# SLIDE 8: UNIT ECONOMICS (metric bars like landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xF8, 0xF8, 0xF8)

section_label(s, "07 / ЭКОНОМИКА")
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "ЮНИТ-ЭКОНОМИКА", size=42, color=BLACK, bold=True, font=HEADING)

# KPI cards — dark cells for bright colors on light bg
kpi_cell(s, Inches(0.8), Inches(2.3), "~120₽", "CAC\n(ОРГАНИКА + ASA)", RED, 200)
kpi_cell_dark(s, Inches(3.7), Inches(2.3), "~1 400₽", "LTV\n(12 МЕС RETENTION)", GREEN, 400)
kpi_cell_dark(s, Inches(6.6), Inches(2.3), "11.7x", "LTV / CAC\nЗДОРОВАЯ ЭКОНОМИКА", GREEN, 600)
kpi_cell_dark(s, Inches(9.5), Inches(2.3), "3%+", "КОНВЕРСИЯ\nFREE → PREMIUM", YELLOW, 800)

# Metric bars (like .metric-bar on landing)
progress_bar(s, Inches(0.8), Inches(4.7), Inches(10.5), 0.40, RGBColor(0x00, 0xAA, 0x44),
             "RETENTION M1", "40%", 1000)
progress_bar(s, Inches(0.8), Inches(5.5), Inches(10.5), 0.15, RGBColor(0xCC, 0xAA, 0x00),
             "RETENTION M12", "15%", 1200)
progress_bar(s, Inches(0.8), Inches(6.3), Inches(10.5), 0.85, RED,
             "GROSS MARGIN", "85%", 1400)


# ────────────────────────────────────────────────────────────
# SLIDE 9: ROADMAP (timeline like landing)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0x0A, 0x0A, 0x0A)

section_label(s, "08 / РОАДМАП", bg_color=RED)
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "27+ СЕССИЙ УЖЕ ПОЗАДИ", size=42, color=WHITE, bold=True, font=HEADING)

phases = [
    ("ГОТОВО", "Core + Multi-trip", "Сессии 1–12", GREEN),
    ("ГОТОВО", "Advanced + Sync + AI", "Сессии 13–27", GREEN),
    ("ДАЛЕЕ", "App Store Launch", "Мес 1–2", RED),
    ("РОСТ", "Growth + Retention", "Мес 3–6", YELLOW),
    ("МАСШТАБ", "Scale + B2B + EN", "Год 2–3", BLUE),
]

# Timeline vertical line
timeline_line = add_rect(s, Inches(1.55), Inches(2.4), Inches(0.06), Inches(4.7),
                         RGBColor(0x33, 0x33, 0x33))
timeline_line.line.fill.background()

for i, (status, phase_title, period, clr) in enumerate(phases):
    y = Inches(2.4 + i * 0.92)

    # Dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.4), y + Inches(0.15),
                              Inches(0.28), Inches(0.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.color.rgb = BLACK
    dot.line.width = Pt(3)

    # Status label
    status_bg = add_rect(s, Inches(2.1), y + Inches(0.12), Inches(1.2), Inches(0.3),
                         clr)
    status_bg.line.fill.background()
    add_text(s, Inches(2.15), y + Inches(0.12), Inches(1.1), Inches(0.3),
             status, size=9, color=BLACK, bold=True, font=MONO, align=PP_ALIGN.CENTER)

    # Title
    add_text(s, Inches(3.6), y + Inches(0.1), Inches(5), Inches(0.4),
             phase_title, size=18, color=WHITE, bold=True, font=HEADING)

    # Period
    add_text(s, Inches(9.5), y + Inches(0.15), Inches(3), Inches(0.3),
             period, size=11, color=GRAY, font=MONO)

    add_entrance_anim(s, dot, 200 + i * 200)
    add_entrance_anim(s, status_bg, 250 + i * 200)


# ────────────────────────────────────────────────────────────
# SLIDE 10: ASK — Investment Round
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xF8, 0xF8, 0xF8)

section_label(s, "09 / РАУНД")
title = add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
                 "ЗАПРАШИВАЕМЫЕ ИНВЕСТИЦИИ", size=42, color=BLACK, bold=True, font=HEADING)

# Left card — Ask amount (brutal card style)
ask_card, ask_shadow = brutal_card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.2), RED)
add_text(s, Inches(1.2), Inches(2.7), Inches(4.7), Inches(0.8),
         "₽2 000 000", size=42, color=RED, bold=True, font=HEADING, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(3.5), Inches(4.7), Inches(0.4),
         "PRE-SEED РАУНД", size=10, color=GRAY, bold=True, font=MONO, align=PP_ALIGN.CENTER)

# Divider inside card
ask_div = add_line(s, Inches(1.5), Inches(4.1), Inches(4), Pt(2), LIGHT_GRAY, thickness=2)

ask_items = [
    "App Store запуск + скриншоты",
    "Apple Search Ads (6 мес)",
    "Серверы (Supabase Pro)",
    "Дизайн + UX-аудит",
    "Юр. лицо + Privacy Policy",
]
ask_lines = add_multiline_text(s, Inches(1.5), Inches(4.3), Inches(4.2), Inches(2.5),
                               [f"•  {item}" for item in ask_items],
                               size=11, color=BODY_TEXT, font=MONO)

add_entrance_anim(s, ask_shadow, 200, 600)
add_entrance_anim(s, ask_card, 200, 600)

# Right card — Goals (dark card for green contrast)
goal_card = add_rect(s, Inches(7), Inches(2.5), Inches(5.5), Inches(4.2),
                      BLACK, GREEN, border_width=3)
add_text(s, Inches(7.4), Inches(2.7), Inches(4.7), Inches(0.6),
         "ЦЕЛИ", size=28, color=GREEN, bold=True, font=HEADING, align=PP_ALIGN.CENTER)

goal_div = add_line(s, Inches(7.7), Inches(3.4), Inches(4), Pt(2), RGBColor(0x33, 0x33, 0x33), thickness=2)

goals = [
    ("30K", "загрузок за 6 мес"),
    ("900", "платящих (3% конверсия)"),
    ("₽315K", "MRR к месяцу 6"),
    ("₽3.8M", "ARR к году 1"),
    ("10x+", "LTV/CAC"),
]
for j, (val, desc) in enumerate(goals):
    gy = Inches(3.65 + j * 0.52)
    add_text(s, Inches(7.5), gy, Inches(1.4), Inches(0.4),
             val, size=16, color=GREEN, bold=True, font=HEADING)
    add_text(s, Inches(8.9), gy + Inches(0.03), Inches(3), Inches(0.4),
             desc, size=11, color=GRAY, font=MONO)

add_entrance_anim(s, goal_card, 600, 600)


# ────────────────────────────────────────────────────────────
# SLIDE 11: CTA (like #cta on landing — red background)
# ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, 0xFF, 0x2D, 0x2D)

# ULTIMA title centered
t1 = add_text(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
              "ULTIMA", size=80, color=WHITE, bold=True, font=HEADING, align=PP_ALIGN.CENTER)
t1b = add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.0),
               "THULE", size=80, color=BLACK, bold=True, font=HEADING, align=PP_ALIGN.CENTER)

t2 = add_text(s, Inches(2), Inches(4.3), Inches(9.3), Inches(0.8),
              "Давайте строить будущее путешествий вместе",
              size=20, color=WHITE, font=MONO, align=PP_ALIGN.CENTER)

# Contact line
contact_line = add_line(s, Inches(4), Inches(5.5), Inches(5.3), Pt(2), WHITE, thickness=2)
t3 = add_text(s, Inches(2), Inches(5.7), Inches(9.3), Inches(0.5),
              "t.me/egorgalkin  •  github.com/russiankendricklamar",
              size=12, color=WHITE, bold=True, font=MONO, align=PP_ALIGN.CENTER)

# Button-like CTA
cta_btn = add_rect(s, Inches(4.5), Inches(6.4), Inches(4.3), Inches(0.6),
                    WHITE, BLACK, border_width=3)
add_text(s, Inches(4.6), Inches(6.42), Inches(4.1), Inches(0.55),
         "СВЯЗАТЬСЯ С НАМИ", size=11, color=BLACK, bold=True, font=MONO, align=PP_ALIGN.CENTER)

add_entrance_anim(s, t1, 0, 600)
add_entrance_anim(s, t1b, 200, 600)
add_entrance_anim(s, t2, 500)
add_entrance_anim(s, t3, 700)
add_entrance_anim(s, cta_btn, 900)


# ════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════
out_path = "/Users/egorgalkin/Travel app/docs/ULTIMA_Pitch_Deck.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
